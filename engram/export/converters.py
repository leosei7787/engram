"""
engram.export.converters — Convert markdown text to common output formats
=========================================================================

Supported formats: md, docx, pptx, pdf (weasyprint), html (pdf fallback)

Usage:
    from engram.export.converters import export
    result = export(text, fmt="pptx", filename="VW_brief", output_dir=Path("outputs/"))
    # → {"path": "/…/VW_brief_20260509_143200.pptx", "name": "...", "format": "pptx"}
"""

from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path
from typing import Literal

ExportFormat = Literal["md", "docx", "pptx", "pdf", "html"]


# ─── Markdown parsing ─────────────────────────────────────────────────────────

def _strip_bold(text: str) -> str:
    return re.sub(r"\*\*(.+?)\*\*", r"\1", text)


def _html_escape(s: str) -> str:
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _html_inline(s: str) -> str:
    s = _html_escape(s)
    s = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", s)
    s = re.sub(r"\*(.+?)\*",     r"<em>\1</em>",         s)
    s = re.sub(r"`(.+?)`",       r"<code>\1</code>",      s)
    return s


def _parse_sections(text: str, doc_title: str = "") -> list[dict]:
    """
    Split markdown into slide-ready sections.
    Each section: {"title": str, "bullets": [str], "is_title": bool, "subtitle": str}
    """
    lines   = text.splitlines()
    sections: list[dict] = []
    current: dict | None = None

    def _flush():
        nonlocal current
        if current is not None:
            sections.append(current)
        current = None

    for raw in lines:
        line = raw.rstrip()
        if not line:
            continue

        # H1 → title slide
        if line.startswith("# "):
            _flush()
            current = {"title": line[2:].strip(), "subtitle": "", "bullets": [], "is_title": True}

        # H2 → content slide
        elif line.startswith("## "):
            _flush()
            current = {"title": line[3:].strip(), "subtitle": "", "bullets": [], "is_title": False}

        # H3 → sub-section label (indented bullet on current slide)
        elif line.startswith("### "):
            if current is None:
                current = {"title": "Summary", "subtitle": "", "bullets": [], "is_title": False}
            current["bullets"].append("  " + line[4:].strip())

        # Bold heading pattern: **Foo:** or **Foo**
        elif re.match(r"^\*\*[^*]+\*\*:?\s*$", line):
            _flush()
            title = line.strip("* :").strip()
            current = {"title": title, "subtitle": "", "bullets": [], "is_title": False}

        # Bullet
        elif line.startswith("- ") or line.startswith("• "):
            if current is None:
                current = {"title": doc_title or "Summary", "subtitle": "", "bullets": [], "is_title": False}
            current["bullets"].append(_strip_bold(line[2:].strip()))

        # Regular text
        else:
            if current is None:
                current = {"title": doc_title or "Summary", "subtitle": "", "bullets": [], "is_title": False}
            clean = _strip_bold(line).strip()
            if not clean:
                continue
            if current.get("is_title") and not current["subtitle"]:
                current["subtitle"] = clean
            else:
                current["bullets"].append(clean)

    _flush()

    if not sections:
        bullets = [_strip_bold(l).strip() for l in lines if l.strip()]
        sections = [{"title": doc_title or "Summary", "subtitle": "", "bullets": bullets, "is_title": False}]

    return sections


# ─── Filename helper ──────────────────────────────────────────────────────────

def _safe_name(name: str, ext: str) -> str:
    ts   = datetime.now().strftime("%Y%m%d_%H%M%S")
    base = re.sub(r"[^a-zA-Z0-9_\- ]", "", name).strip().replace(" ", "_")[:48]
    return f"{base}_{ts}.{ext}" if base else f"output_{ts}.{ext}"


# ─── Markdown ─────────────────────────────────────────────────────────────────

def to_md(text: str, filename: str, output_dir: Path) -> Path:
    out = output_dir / _safe_name(filename, "md")
    out.write_text(text, encoding="utf-8")
    return out


# ─── DOCX ─────────────────────────────────────────────────────────────────────

def to_docx(text: str, filename: str, output_dir: Path) -> Path:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor

    doc = Document()

    # Margins
    for sec in doc.sections:
        sec.top_margin    = Inches(1.0)
        sec.bottom_margin = Inches(1.0)
        sec.left_margin   = Inches(1.2)
        sec.right_margin  = Inches(1.2)

    for raw in text.splitlines():
        line = raw.rstrip()
        if not line:
            continue
        if line.startswith("# "):
            doc.add_heading(_strip_bold(line[2:]), level=0)
        elif line.startswith("## "):
            doc.add_heading(_strip_bold(line[3:]), level=1)
        elif line.startswith("### "):
            doc.add_heading(_strip_bold(line[4:]), level=2)
        elif line.startswith("- ") or line.startswith("• "):
            p = doc.add_paragraph(style="List Bullet")
            _docx_add_runs(p, line[2:])
        else:
            p = doc.add_paragraph()
            _docx_add_runs(p, line)

    out = output_dir / _safe_name(filename, "docx")
    doc.save(out)
    return out


def _docx_add_runs(para, text: str) -> None:
    """Split on **bold** markers and add formatted runs."""
    for i, part in enumerate(re.split(r"\*\*(.+?)\*\*", text)):
        if not part:
            continue
        run = para.add_run(part)
        run.bold = (i % 2 == 1)


# ─── PPTX ─────────────────────────────────────────────────────────────────────

def to_pptx(text: str, filename: str, output_dir: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    sections = _parse_sections(text, filename.replace("_", " "))

    # Ensure at least one title slide
    if not sections[0].get("is_title"):
        sections.insert(0, {
            "title":    filename.replace("_", " ").title(),
            "subtitle": datetime.now().strftime("Generated %d %b %Y"),
            "bullets":  [],
            "is_title": True,
        })

    for sec in sections:
        if sec.get("is_title"):
            sl = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
            sl.shapes.title.text = sec["title"]
            try:
                sl.placeholders[1].text = sec.get("subtitle", "")
            except (KeyError, IndexError):
                pass
        else:
            sl = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
            sl.shapes.title.text = sec["title"]
            try:
                tf = sl.placeholders[1].text_frame
                tf.clear()
                for j, bullet in enumerate(sec.get("bullets", [])[:14]):
                    p       = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    is_sub  = bullet.startswith("  ")
                    p.text  = bullet.strip()
                    p.level = 1 if is_sub else 0
            except (KeyError, IndexError):
                pass

    out = output_dir / _safe_name(filename, "pptx")
    prs.save(out)
    return out


# ─── HTML (styled, printable to PDF) ─────────────────────────────────────────

def to_html(text: str, title: str) -> str:
    parts: list[str] = []
    in_list = False

    for raw in text.splitlines():
        line = raw.rstrip()
        is_bullet = line.startswith("- ") or line.startswith("• ")

        if in_list and not is_bullet:
            parts.append("</ul>")
            in_list = False

        if not line:
            continue

        if line.startswith("# "):
            parts.append(f"<h1>{_html_inline(line[2:])}</h1>")
        elif line.startswith("## "):
            parts.append(f"<h2>{_html_inline(line[3:])}</h2>")
        elif line.startswith("### "):
            parts.append(f"<h3>{_html_inline(line[4:])}</h3>")
        elif is_bullet:
            if not in_list:
                parts.append("<ul>")
                in_list = True
            parts.append(f"<li>{_html_inline(line[2:])}</li>")
        else:
            parts.append(f"<p>{_html_inline(line)}</p>")

    if in_list:
        parts.append("</ul>")

    body = "\n".join(parts)
    return f"""<!DOCTYPE html>
<html lang="en"><head>
<meta charset="UTF-8">
<title>{_html_escape(title)}</title>
<style>
  body {{ font-family: -apple-system, Arial, sans-serif; max-width: 820px;
          margin: 48px auto; padding: 0 28px; color: #1a1a1a; line-height: 1.65; }}
  h1 {{ font-size: 28px; font-weight: 800; border-bottom: 3px solid #D97757;
        padding-bottom: 10px; margin-bottom: 20px; }}
  h2 {{ font-size: 18px; font-weight: 700; color: #D97757; margin-top: 32px; margin-bottom: 8px; }}
  h3 {{ font-size: 15px; font-weight: 600; color: #555; margin-bottom: 6px; }}
  p  {{ margin: 6px 0; }}
  ul {{ padding-left: 22px; margin: 8px 0; }}
  li {{ margin-bottom: 5px; }}
  strong {{ font-weight: 700; }}
  code {{ background: #f3f3f3; padding: 1px 5px; border-radius: 3px;
           font-family: monospace; font-size: 13px; }}
  @media print {{ body {{ margin: 0; }} }}
</style>
</head>
<body>
{body}
</body></html>"""


def to_html_file(text: str, filename: str, output_dir: Path) -> Path:
    html  = to_html(text, filename)
    out   = output_dir / _safe_name(filename, "html")
    out.write_text(html, encoding="utf-8")
    return out


def to_pdf(text: str, filename: str, output_dir: Path) -> Path | None:
    """PDF via weasyprint; returns None if not installed."""
    try:
        import weasyprint
        html = to_html(text, filename)
        out  = output_dir / _safe_name(filename, "pdf")
        weasyprint.HTML(string=html).write_pdf(str(out))
        return out
    except ImportError:
        return None


# ─── Main export entry point ──────────────────────────────────────────────────

def export(text: str, fmt: str, filename: str, output_dir: Path) -> dict:
    """
    Export text to the given format, save to output_dir, return metadata dict.

    Returns:
        {"path": str, "name": str, "format": str}          on success
        {"path": ...,  "note": str, "format": "html"}       when PDF falls back to HTML
        {"error": str}                                       on failure
    """
    output_dir.mkdir(parents=True, exist_ok=True)

    try:
        if fmt == "md":
            path = to_md(text, filename, output_dir)
            return {"path": str(path), "name": path.name, "format": "md"}

        elif fmt == "docx":
            path = to_docx(text, filename, output_dir)
            return {"path": str(path), "name": path.name, "format": "docx"}

        elif fmt == "pptx":
            path = to_pptx(text, filename, output_dir)
            return {"path": str(path), "name": path.name, "format": "pptx"}

        elif fmt == "pdf":
            path = to_pdf(text, filename, output_dir)
            if path:
                return {"path": str(path), "name": path.name, "format": "pdf"}
            # Fallback to HTML (browser can print to PDF)
            path = to_html_file(text, filename, output_dir)
            return {
                "path":   str(path),
                "name":   path.name,
                "format": "html",
                "note":   "weasyprint not installed — saved as HTML. Open and use File → Print → Save as PDF.",
            }

        elif fmt == "html":
            path = to_html_file(text, filename, output_dir)
            return {"path": str(path), "name": path.name, "format": "html"}

        else:
            return {"error": f"Unknown format: {fmt!r}"}

    except Exception as e:
        return {"error": str(e)}

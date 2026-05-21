"""
engram.ingest.extractors — text extraction for Office docs and images
======================================================================

The watcher's ``_read_file`` originally handled only ``.pdf`` (via ``pypdf``)
plus any text-shaped file via raw read. Everything else — Word docs,
PowerPoints, Excel sheets, screenshots — landed in the inbox and was
silently skipped by the extension allowlist. This module fills the gap.

Each extractor returns a single plain-text string. Empty string means the
extractor ran but found no text (which the caller treats as "skip, no
signal"); raising means a real error to log.

  - .docx → ``python-docx``     (paragraphs + simple table cells)
  - .pptx → ``python-pptx``     (each slide as a section, text frames + tables + notes)
  - .xlsx → ``openpyxl``        (each sheet as a table; numbers + text)
  - images → Anthropic vision  (description + extracted text via the SDK,
              skipped with a warning if ``ANTHROPIC_API_KEY`` is unset —
              we deliberately do NOT shell out to tesseract because it
              isn't installed on this box and OCR quality on exec
              screenshots is poor anyway)

These libraries are already declared in ``requirements.txt`` (they were
added for the EXPORT pipeline — we reuse them on the INGEST side here).
"""
from __future__ import annotations

import base64
import os
import traceback
from pathlib import Path
from typing import Optional


# ─── .docx ────────────────────────────────────────────────────────────────────

def extract_docx(path: Path) -> str:
    """Extract paragraphs + table-cell text from a .docx file."""
    try:
        import docx as _docx
    except ImportError:
        return ""
    try:
        doc = _docx.Document(str(path))
    except Exception:
        print(f"[extractors] docx open failed for {path.name}:\n" + traceback.format_exc(), flush=True)
        return ""
    parts: list[str] = []
    for para in doc.paragraphs:
        t = (para.text or "").strip()
        if t:
            parts.append(t)
    # Tables — flatten as "cell | cell | cell" rows. Keeps tabular data
    # searchable without trying to reproduce layout.
    for tbl in doc.tables:
        for row in tbl.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                parts.append(line)
    return "\n\n".join(parts)


# ─── .pptx ────────────────────────────────────────────────────────────────────

def extract_pptx(path: Path) -> str:
    """Extract slide text + speaker notes from a .pptx file.

    Each slide becomes a section so the slide structure is preserved in the
    output (lets the curator scan recognise which content came from which slide).
    """
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    try:
        prs = Presentation(str(path))
    except Exception:
        print(f"[extractors] pptx open failed for {path.name}:\n" + traceback.format_exc(), flush=True)
        return ""
    sections: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_bits: list[str] = []
        # Walk every shape and pull any text it owns. ``shape.has_text_frame``
        # catches title + body + arbitrary text boxes; tables we handle below.
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame:
                t = (shape.text_frame.text or "").strip()
                if t:
                    slide_bits.append(t)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [(c.text or "").strip() for c in row.cells]
                    line = " | ".join(c for c in cells if c)
                    if line:
                        slide_bits.append(line)
        # Speaker notes — often where the actual reasoning lives in exec decks
        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes = ""
        if not slide_bits and not notes:
            continue
        section = f"## Slide {i}\n\n" + "\n\n".join(slide_bits)
        if notes:
            section += f"\n\n_Speaker notes:_ {notes}"
        sections.append(section)
    return "\n\n---\n\n".join(sections)


# ─── .xlsx ────────────────────────────────────────────────────────────────────

def extract_xlsx(path: Path, *, max_rows_per_sheet: int = 500) -> str:
    """Extract cells from every sheet of an .xlsx. Caps rows per sheet to keep
    chat-context budgets sane — most exec spreadsheets are <500 rows, and
    bigger ones tend to be raw data dumps where row 500+ rarely adds signal.
    """
    try:
        import openpyxl
    except ImportError:
        return ""
    try:
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    except Exception:
        print(f"[extractors] xlsx open failed for {path.name}:\n" + traceback.format_exc(), flush=True)
        return ""
    sections: list[str] = []
    for sheet_name in wb.sheetnames:
        try:
            ws = wb[sheet_name]
        except Exception:
            continue
        rows_out: list[str] = []
        for r, row in enumerate(ws.iter_rows(values_only=True)):
            if r >= max_rows_per_sheet:
                rows_out.append(f"_(truncated at {max_rows_per_sheet} rows)_")
                break
            line = " | ".join("" if v is None else str(v) for v in row)
            if line.replace("|", "").strip():
                rows_out.append(line)
        if rows_out:
            sections.append(f"## Sheet: {sheet_name}\n\n" + "\n".join(rows_out))
    try:
        wb.close()
    except Exception:
        pass
    return "\n\n---\n\n".join(sections)


# ─── images (vision) ─────────────────────────────────────────────────────────

_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".gif", ".bmp", ".webp"}


def _media_type_for(suffix: str) -> str:
    s = suffix.lower()
    if s in (".jpg", ".jpeg"): return "image/jpeg"
    if s == ".png":            return "image/png"
    if s == ".gif":            return "image/gif"
    if s == ".webp":           return "image/webp"
    if s in (".tiff", ".tif"): return "image/tiff"
    if s == ".bmp":            return "image/bmp"
    return "application/octet-stream"


def extract_image(path: Path, *, model: str = "claude-haiku-4-5",
                  max_bytes: int = 5_000_000) -> str:
    """Run the Anthropic vision SDK over an image and return a markdown summary.

    Returns an empty string when ``ANTHROPIC_API_KEY`` is not set (no SDK
    fallback available without API access on this box — tesseract isn't
    installed). The caller treats empty as "skipped, no signal" and writes
    nothing to MEMORY.

    For business-document screenshots this gives surprisingly searchable
    output: text content gets transcribed, charts get described, plus a
    short topical caption. The keyword scanner can then surface the image
    when relevant queries hit.
    """
    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        print(f"[extractors] image skipped (no ANTHROPIC_API_KEY): {path.name}", flush=True)
        return ""
    try:
        size = path.stat().st_size
    except Exception:
        return ""
    if size > max_bytes:
        print(f"[extractors] image too large ({size} bytes), skipping: {path.name}", flush=True)
        return ""

    try:
        import anthropic
    except ImportError:
        return ""

    try:
        data = path.read_bytes()
        b64 = base64.standard_b64encode(data).decode("ascii")
    except Exception:
        return ""

    media_type = _media_type_for(path.suffix)
    prompt = (
        "You are extracting durable signal from an image attached to a knowledge memory store. "
        "The image is most likely an exec screenshot, slide capture, dashboard, chart, or "
        "scanned doc. Output strict markdown only (no preamble, no closing pleasantry):\n\n"
        "## Caption\n<one-sentence description of what the image is>\n\n"
        "## Extracted text\n<every legible word/line, transcribed verbatim, line-by-line; "
        "omit boilerplate UI chrome>\n\n"
        "## Data / entities\n<bullets: any numbers, dates, named people, projects, "
        "accounts, or named programs visible. Skip if none.>\n\n"
        "## Topic\n<3-7 keyword tags, comma-separated — used by a keyword retrieval scanner>"
    )

    try:
        client = anthropic.Anthropic(api_key=api_key)
        resp = client.messages.create(
            model=model,
            max_tokens=1500,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "image", "source": {
                        "type": "base64", "media_type": media_type, "data": b64,
                    }},
                    {"type": "text", "text": prompt},
                ],
            }],
        )
    except Exception:
        print(f"[extractors] vision call failed for {path.name}:\n" + traceback.format_exc(), flush=True)
        return ""

    out = ""
    for block in (resp.content or []):
        if getattr(block, "type", "") == "text":
            out += getattr(block, "text", "")
    return out.strip()


# ─── Dispatch ─────────────────────────────────────────────────────────────────

def extract(path: Path) -> str:
    """Pick the right extractor by extension. Returns text (possibly empty).

    Unhandled extensions return an empty string — caller decides what to do
    (typically: skip; treat as a placeholder).
    """
    ext = path.suffix.lower()
    if ext == ".docx": return extract_docx(path)
    if ext == ".pptx": return extract_pptx(path)
    if ext == ".xlsx": return extract_xlsx(path)
    if ext in _IMAGE_EXTENSIONS: return extract_image(path)
    return ""


# Public list — used by the watcher to set its extensions allowlist
SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".xlsx", *_IMAGE_EXTENSIONS}
